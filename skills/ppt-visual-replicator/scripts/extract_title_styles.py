#!/usr/bin/env python3
"""Extract source-title typography so reconstruction cannot invent it."""

from __future__ import annotations

import argparse
import json
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE


def _rgb_hex(font) -> str | None:
    try:
        if font.color.rgb:
            return str(font.color.rgb)
    except Exception:
        pass
    xml = getattr(font, "_element", None)
    xml = getattr(xml, "xml", "")
    for color in ("srgbClr", "prstClr"):
        marker = f'<a:{color} val="'
        start = xml.find(marker)
        if start >= 0:
            start += len(marker)
            end = xml.find('"', start)
            return xml[start:end].upper()
    if 'schemeClr val="tx1"' in xml or 'schemeClr val="dk1"' in xml:
        return "000000"
    if 'schemeClr val="bg1"' in xml or 'schemeClr val="lt1"' in xml or 'prstClr val="white"' in xml:
        return "FFFFFF"
    return None


def _walk_shapes(shapes):
    for shape in shapes:
        try:
            shape_type = shape.shape_type
        except NotImplementedError:
            shape_type = None
        if shape_type == MSO_SHAPE_TYPE.GROUP:
            yield from _walk_shapes(shape.shapes)
        else:
            yield shape


def _text_shape_style(shape, slide_width: int, slide_height: int) -> dict[str, Any] | None:
    if not getattr(shape, "has_text_frame", False) or not shape.text.strip():
        return None
    runs = [run for paragraph in shape.text_frame.paragraphs for run in paragraph.runs if run.text]
    if not runs:
        return None
    dominant = max(runs, key=lambda run: float(run.font.size.pt) if run.font.size else 0.0)
    font_size = float(dominant.font.size.pt) if dominant.font.size else 0.0
    return {
        "text": shape.text.strip(),
        "font": dominant.font.name,
        "font_size_pt": font_size,
        "color": _rgb_hex(dominant.font),
        "bold": bool(dominant.font.bold),
        "italic": bool(dominant.font.italic),
        "box_ratio": {
            "left": round(shape.left / slide_width, 5),
            "top": round(shape.top / slide_height, 5),
            "width": round(shape.width / slide_width, 5),
            "height": round(shape.height / slide_height, 5),
        },
    }


def extract_title_styles(target: str | Path, target_slide: int | None = None) -> dict[str, Any]:
    source = Path(target).expanduser().resolve()
    try:
        presentation = Presentation(source)
    except (KeyError, OSError, ValueError) as exc:
        # PPTX inspection can still be sufficient for direct-run planning on
        # minimal OOXML fixtures. A missing high-fidelity title record must not
        # prevent that established workflow; real PowerPoint files continue to
        # provide the source typography record below.
        return {
            "schema": "ppt_visual_title_styles.v1",
            "target_pptx": str(source),
            "slides": [],
            "unavailable_reason": f"native title style extraction unavailable: {exc}",
            "rule": "Use the declared source title style unless an explicit style-contract title_system overrides that title role.",
        }
    if target_slide is not None and not 1 <= target_slide <= len(presentation.slides):
        raise ValueError(f"target slide does not exist: {target_slide}")
    records = []
    for number, slide in enumerate(presentation.slides, start=1):
        if target_slide is not None and number != target_slide:
            continue
        candidates = []
        for shape in _walk_shapes(slide.shapes):
            style = _text_shape_style(shape, presentation.slide_width, presentation.slide_height)
            if not style:
                continue
            top = style["box_ratio"]["top"]
            if top > 0.35:
                continue
            # A page title is ordinarily the largest text in the upper third.
            # Penalize short section ribbons only after preserving their style
            # as secondary candidates for worker inspection.
            score = style["font_size_pt"] * 10 + (1 - top) * 4 + min(len(style["text"]), 120) / 120
            candidates.append((score, style))
        candidates.sort(key=lambda item: item[0], reverse=True)
        records.append(
            {
                "slide_number": number,
                "title": candidates[0][1] if candidates else None,
                "secondary_candidates": [item[1] for item in candidates[1:4]],
            }
        )
    return {
        "schema": "ppt_visual_title_styles.v1",
        "target_pptx": str(source),
        "slides": records,
        "rule": "Use the declared source title style unless an explicit style-contract title_system overrides that title role.",
    }


def main() -> int:
    parser = argparse.ArgumentParser(description="Extract native title font, color, size, weight, and placement from a PPTX.")
    parser.add_argument("--target", required=True)
    parser.add_argument("--target-slide", type=int)
    parser.add_argument("--out", required=True)
    args = parser.parse_args()
    try:
        payload = extract_title_styles(args.target, args.target_slide)
    except (OSError, ValueError) as exc:
        raise SystemExit(str(exc)) from exc
    output = Path(args.out).expanduser().resolve()
    output.parent.mkdir(parents=True, exist_ok=True)
    output.write_text(json.dumps(payload, ensure_ascii=False, indent=2) + "\n", encoding="utf-8")
    print(json.dumps(payload, ensure_ascii=False, indent=2))
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
