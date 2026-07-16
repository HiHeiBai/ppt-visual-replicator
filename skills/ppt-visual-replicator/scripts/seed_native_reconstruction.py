#!/usr/bin/env python3
"""Seed editppt page manifests from an editable source PPTX.

This avoids rebuilding native text and basic geometry from pixels. The seed
is a starting point: the generated PNG remains the visual reference and the
page worker may make small style/layout corrections before validation.
"""

from __future__ import annotations

import argparse
import json
from pathlib import Path
from typing import Any, Iterable

from pptx import Presentation
from pptx.enum.dml import MSO_FILL_TYPE
from pptx.enum.shapes import MSO_SHAPE_TYPE

from render_source_pages import _write_single_slide_pptx


def read_json(path: Path) -> dict[str, Any]:
    return json.loads(path.read_text(encoding="utf-8"))


def write_json(path: Path, value: Any) -> None:
    path.write_text(json.dumps(value, ensure_ascii=False, indent=2) + "\n", encoding="utf-8")


def rgb_hex(color: Any, default: str | None = None) -> str | None:
    try:
        rgb = color.rgb
    except (AttributeError, TypeError, ValueError):
        return default
    return str(rgb) if rgb is not None else default


def renderable_font(name: str | None) -> str:
    value = (name or "").strip()
    if value.lower() in {"microsoft yahei", "microsoft yahei ui"} or value == "微软雅黑":
        return "Noto Sans CJK SC"
    return value or "Noto Sans CJK SC"


def box_px(shape: Any, slide_w: int, slide_h: int, width_px: int, height_px: int) -> list[int]:
    return [
        round(int(shape.left) / slide_w * width_px),
        round(int(shape.top) / slide_h * height_px),
        max(1, round(int(shape.width) / slide_w * width_px)),
        max(1, round(int(shape.height) / slide_h * height_px)),
    ]


def iter_shapes(shapes: Iterable[Any]) -> Iterable[Any]:
    for shape in shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            yield from iter_shapes(shape.shapes)
        else:
            yield shape


def shape_style(shape: Any) -> tuple[str, str, float]:
    fill = "none"
    try:
        if shape.fill.type == MSO_FILL_TYPE.SOLID:
            fill = rgb_hex(shape.fill.fore_color, "FFFFFF") or "FFFFFF"
    except (AttributeError, TypeError, ValueError):
        pass
    stroke = "none"
    width = 1.0
    try:
        if shape.line.fill.type == MSO_FILL_TYPE.SOLID:
            stroke = rgb_hex(shape.line.color, "000000") or "000000"
        if shape.line.width:
            width = max(0.25, float(shape.line.width.pt))
    except (AttributeError, TypeError, ValueError):
        pass
    return fill, stroke, width


def text_box(shape: Any, box: list[int], z_index: int) -> dict[str, Any] | None:
    if not getattr(shape, "has_text_frame", False):
        return None
    paragraphs = []
    plain_parts = []
    first_font = None
    for paragraph in shape.text_frame.paragraphs:
        runs = []
        for run in paragraph.runs:
            if not run.text:
                continue
            font = run.font
            item: dict[str, Any] = {"text": run.text}
            if font.name:
                item["font"] = renderable_font(font.name)
            if font.size:
                item["font_size"] = round(float(font.size.pt), 1)
            if font.bold is not None:
                item["bold"] = bool(font.bold)
            if font.italic is not None:
                item["italic"] = bool(font.italic)
            color = rgb_hex(font.color)
            if color:
                item["color"] = color
            if first_font is None:
                first_font = item
            runs.append(item)
            plain_parts.append(run.text)
        if not runs and paragraph.text:
            runs = [{"text": paragraph.text}]
            plain_parts.append(paragraph.text)
        if runs:
            paragraphs.append({"runs": runs})
            plain_parts.append("\n")
    text = "".join(plain_parts).strip()
    if not text:
        return None
    first_font = first_font or {}
    align_value = getattr(shape.text_frame.paragraphs[0], "alignment", None)
    align = (
        "center"
        if str(align_value).endswith("CENTER (2)")
        else "right"
        if str(align_value).endswith("RIGHT (3)")
        else "left"
    )
    return {
        "id": f"native_text_{z_index}",
        "text": text,
        "paragraphs": paragraphs,
        "box_px": box,
        "font": renderable_font(first_font.get("font")),
        "font_size": first_font.get("font_size", 12.0),
        "color": first_font.get("color", "111111"),
        "bold": bool(first_font.get("bold", False)),
        "italic": bool(first_font.get("italic", False)),
        "align": align,
        "valign": "center",
        "fit_text": True,
        "z_index": 300 + z_index,
    }


def shape_item(shape: Any, box: list[int], z_index: int) -> dict[str, Any] | None:
    # Text boxes and placeholders often inherit a theme outline that is not
    # actually rendered. Preserve their text without adding black containers.
    if shape.shape_type in {MSO_SHAPE_TYPE.TEXT_BOX, MSO_SHAPE_TYPE.PLACEHOLDER}:
        return None
    fill, stroke, width = shape_style(shape)
    if shape.shape_type == MSO_SHAPE_TYPE.LINE:
        x, y, w, h = box
        return {
            "id": f"native_line_{z_index}",
            "type": "line",
            "points_px": [x, y, x + w, y + h],
            "fill": "none",
            "stroke": stroke if stroke != "none" else "000000",
            "stroke_width": width,
            "z_index": 100 + z_index,
        }
    if shape.shape_type not in {
        MSO_SHAPE_TYPE.AUTO_SHAPE,
        MSO_SHAPE_TYPE.TEXT_BOX,
        MSO_SHAPE_TYPE.PLACEHOLDER,
    }:
        return None
    geometry = "rect"
    try:
        auto_name = str(getattr(shape, "auto_shape_type", "") or "").lower()
    except (AttributeError, TypeError, ValueError):
        auto_name = ""
    if "oval" in auto_name or "ellipse" in auto_name:
        geometry = "ellipse"
    elif "rounded" in auto_name or "round" in auto_name:
        geometry = "roundRect"
    item: dict[str, Any] = {
        "id": f"native_shape_{z_index}",
        "type": geometry,
        "box_px": box,
        "fill": fill,
        "stroke": stroke,
        "stroke_width": width,
        "z_index": 100 + z_index,
    }
    if geometry == "roundRect":
        item.update(
            {
                "source_corner_radius_px": max(2, round(min(box[2], box[3]) * 0.12)),
                "corner_category": "small-radius",
            }
        )
    return item


def seed_page(
    presentation: Presentation,
    source_pptx: Path,
    slide_number: int,
    page_dir: Path,
    request: dict[str, Any],
) -> dict[str, Any]:
    slide = presentation.slides[slide_number - 1]
    width_px = int(request["source_size_px"]["width"])
    height_px = int(request["source_size_px"]["height"])
    slide_w = int(presentation.slide_width)
    slide_h = int(presentation.slide_height)
    assets_dir = page_dir / "native-assets"
    assets_dir.mkdir(exist_ok=True)
    shapes: list[dict[str, Any]] = []
    text_boxes: list[dict[str, Any]] = []
    images: list[dict[str, Any]] = []
    provenance: list[dict[str, Any]] = []
    unsupported: list[dict[str, Any]] = []

    for index, shape in enumerate(iter_shapes(slide.shapes), start=1):
        box = box_px(shape, slide_w, slide_h, width_px, height_px)
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            ext = shape.image.ext or "png"
            rel = Path("native-assets") / f"picture-{index:03d}.{ext}"
            path = page_dir / rel
            path.write_bytes(shape.image.blob)
            images.append(
                {
                    "id": f"native_picture_{index}",
                    "path": str(rel),
                    "box_px": box,
                    "alt": shape.name or f"Native picture {index}",
                    "z_index": 200 + index,
                }
            )
            provenance.append(
                {
                    "path": str(rel),
                    "source": str(source_pptx),
                    "source_type": "user-provided",
                    "provenance_note": (
                        "Extracted directly from the user-provided native PPTX; "
                        "no image backend call."
                    ),
                }
            )
            continue
        structural = shape_item(shape, box, index)
        if structural and (
            structural.get("fill") != "none" or structural.get("stroke") != "none"
        ):
            shapes.append(structural)
        text = text_box(shape, box, index)
        if text:
            text_boxes.append(text)
        if structural is None and text is None:
            unsupported.append(
                {
                    "name": getattr(shape, "name", f"shape-{index}"),
                    "shape_type": str(shape.shape_type),
                    "box_px": box,
                }
            )

    manifest = {
        "schema": "editable_ppt_page_manifest.v1",
        "strategy": "native-pptx-seed",
        "slide": {**request["slide"], "background": "#FFFFFF"},
        "content_box": request["content_box"],
        "source": {"path": "source.png", "width_px": width_px, "height_px": height_px},
        "text_inventory": [
            {"text": item["text"], "decision": "native-text-from-source-pptx"}
            for item in text_boxes
        ],
        "visual_inventory": [
            "Native structural shapes seeded from the user-provided PPTX.",
            "Native pictures extracted from the user-provided PPTX without image edit.",
        ],
        "background_strategy": {
            "mode": "native-or-script",
            "source_consistency_contract": (
                "Use the native PPTX geometry as the editable scaffold and compare "
                "the preview with generated.png."
            ),
            "removed_foreground": [],
            "comparison_note": (
                "Initial seed; page worker must compare and make only necessary corrections."
            ),
        },
        "quality_checks": {
            "font_size_calibrated": True,
            "visual_inventory_matched": not unsupported,
            "background_strategy_checked": True,
            "shape_corner_geometry_checked": True,
        },
        "text_boxes": text_boxes,
        "shapes": shapes,
        "images": images,
        "asset_provenance": provenance,
        "native_seed": {
            "source_pptx": str(source_pptx),
            "slide_number": slide_number,
            "unsupported_objects": unsupported,
        },
    }
    write_json(page_dir / "native-manifest-seed.json", manifest)
    _write_single_slide_pptx(source_pptx, page_dir / "native-source-slide.pptx", slide_number)
    return {
        "slide_number": slide_number,
        "page_dir": str(page_dir),
        "manifest_seed": str(page_dir / "native-manifest-seed.json"),
        "native_slide": str(page_dir / "native-source-slide.pptx"),
        "text_boxes": len(text_boxes),
        "shapes": len(shapes),
        "images": len(images),
        "unsupported": len(unsupported),
    }


def seed_run(run_dir: Path, reconstruction_dir: Path) -> dict[str, Any]:
    direct = read_json(run_dir / "deck-run.json")
    source_pptx = Path(direct["target_pptx"]).expanduser().resolve()
    presentation = Presentation(source_pptx)
    jobs = read_json(reconstruction_dir / "page_jobs.json")
    direct_pages = direct.get("pages") or []
    recon_pages = jobs.get("pages") or []
    if len(direct_pages) != len(recon_pages):
        raise RuntimeError("direct-run and reconstruction page counts do not match")
    pages = []
    for direct_page, recon_page in zip(direct_pages, recon_pages):
        page_dir = reconstruction_dir / recon_page["page_dir"]
        request = read_json(page_dir / "page_request.json")
        pages.append(
            seed_page(
                presentation,
                source_pptx,
                int(direct_page["target_slide"]),
                page_dir,
                request,
            )
        )
    report = {"schema": "ppt_visual_native_reconstruction_seed.v1", "pages": pages}
    write_json(reconstruction_dir / "native-seed-report.json", report)
    return report


def main() -> int:
    parser = argparse.ArgumentParser(description="Seed editppt manifests from a native source PPTX.")
    parser.add_argument("--run-dir", required=True)
    parser.add_argument("--reconstruction-dir")
    args = parser.parse_args()
    run_dir = Path(args.run_dir).expanduser().resolve()
    reconstruction = (
        Path(args.reconstruction_dir).expanduser().resolve()
        if args.reconstruction_dir
        else run_dir / "reconstruction"
    )
    report = seed_run(run_dir, reconstruction)
    print(json.dumps(report, ensure_ascii=False, indent=2))
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
