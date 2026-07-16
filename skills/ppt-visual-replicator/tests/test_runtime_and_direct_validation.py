from __future__ import annotations

import hashlib
import json
import sys
import tempfile
import unittest
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt


SKILL_ROOT = Path(__file__).resolve().parents[1]
SCRIPTS = SKILL_ROOT / "scripts"
RUNTIME = SKILL_ROOT / "reconstruction" / "cli" / "editppt" / "runtime"
sys.path[:0] = [str(SCRIPTS), str(RUNTIME)]

from extract_title_styles import extract_title_styles  # noqa: E402
from prepare_direct_deck import prepare_direct_deck  # noqa: E402
from record_page_result import validate_direct_visual_gate  # noqa: E402


def sha256(path: Path) -> str:
    return hashlib.sha256(path.read_bytes()).hexdigest()


def make_native_deck(path: Path, titles: list[str]) -> None:
    deck = Presentation()
    for title in titles:
        slide = deck.slides.add_slide(deck.slide_layouts[6])
        box = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(10), Inches(0.6))
        run = box.text_frame.paragraphs[0].add_run()
        run.text = title
        run.font.name = "Microsoft YaHei"
        run.font.size = Pt(27)
        run.font.bold = True
        run.font.color.rgb = RGBColor(0x00, 0x55, 0x87)
    deck.save(path)


class NativePreservationTests(unittest.TestCase):
    def test_extracts_title_source_style(self) -> None:
        with tempfile.TemporaryDirectory() as temp_dir:
            target = Path(temp_dir) / "source.pptx"
            make_native_deck(target, ["MAJESTEC-9：标题样式"])

            payload = extract_title_styles(target)

            title = payload["slides"][0]["title"]
            self.assertEqual(title["font"], "Microsoft YaHei")
            self.assertEqual(title["font_size_pt"], 27.0)
            self.assertEqual(title["color"], "005587")
            self.assertTrue(title["bold"])

    def test_preserves_a_fully_native_deck_without_redraw(self) -> None:
        with tempfile.TemporaryDirectory() as temp_dir:
            root = Path(temp_dir)
            target = root / "source.pptx"
            make_native_deck(target, ["第一页"])

            result = prepare_direct_deck(target, root / "run")
            final = root / "run" / "reconstruction" / "final" / "origin_edited.pptx"

            self.assertEqual(result["status"], "native-source-preserved")
            self.assertTrue(final.is_file())
            self.assertEqual(sha256(final), sha256(target))

    def test_preserves_one_selected_native_slide_as_one_slide_pptx(self) -> None:
        with tempfile.TemporaryDirectory() as temp_dir:
            root = Path(temp_dir)
            target = root / "source.pptx"
            make_native_deck(target, ["第一页", "第二页"])

            result = prepare_direct_deck(target, root / "run", target_slide=2)
            final = root / "run" / "reconstruction" / "final" / "origin_edited.pptx"
            selected = Presentation(final)

            self.assertEqual(result["target_slide_numbers"], [2])
            self.assertEqual(len(selected.slides), 1)
            self.assertIn("第二页", selected.slides[0].shapes[0].text)


class VisualGateTests(unittest.TestCase):
    def test_direct_run_accepts_a_matching_original_source_gate(self) -> None:
        with tempfile.TemporaryDirectory() as temp_dir:
            direct_root = Path(temp_dir) / "run"
            run_dir = direct_root / "reconstruction"
            page_dir = run_dir / "pages" / "page_001"
            source = direct_root / "pages" / "slide-001" / "source-content.png"
            page_pptx = page_dir / "page.pptx"
            source.parent.mkdir(parents=True)
            page_dir.mkdir(parents=True)
            source.write_bytes(b"original-source")
            page_pptx.write_bytes(b"rebuilt-page")
            (direct_root / "deck-run.json").write_text(
                json.dumps({"pages": [{"source_image": "pages/slide-001/source-content.png"}]}),
                encoding="utf-8",
            )
            gate = {
                "schema": "ppt_visual_page_gate.v2",
                "passed": True,
                "manual_accept": True,
                "visual_accept": True,
                "content_accept": True,
                "source_sha256": sha256(source),
                "content_source_sha256": sha256(source),
                "page_pptx_sha256": sha256(page_pptx),
                "checks": {
                    "visual_reference_match": True,
                    "original_content_reviewed": True,
                },
            }
            (page_dir / "visual-gate.json").write_text(json.dumps(gate), encoding="utf-8")

            evidence = validate_direct_visual_gate(
                run_dir, {"page_id": "page_001", "page_index": 1}, page_dir, page_pptx
            )

            self.assertEqual(evidence["path"], "pages/page_001/visual-gate.json")

    def test_direct_run_rejects_a_visual_only_gate(self) -> None:
        with tempfile.TemporaryDirectory() as temp_dir:
            direct_root = Path(temp_dir) / "run"
            run_dir = direct_root / "reconstruction"
            page_dir = run_dir / "pages" / "page_001"
            source = direct_root / "pages" / "slide-001" / "source-content.png"
            page_pptx = page_dir / "page.pptx"
            source.parent.mkdir(parents=True)
            page_dir.mkdir(parents=True)
            source.write_bytes(b"original-source")
            page_pptx.write_bytes(b"rebuilt-page")
            (direct_root / "deck-run.json").write_text(
                json.dumps({"pages": [{"source_image": "pages/slide-001/source-content.png"}]}),
                encoding="utf-8",
            )
            gate = {
                "schema": "ppt_visual_page_gate.v2",
                "passed": True,
                "manual_accept": True,
                "visual_accept": True,
                "content_accept": False,
                "source_sha256": sha256(source),
                "content_source_sha256": sha256(source),
                "page_pptx_sha256": sha256(page_pptx),
                "checks": {
                    "visual_reference_match": True,
                    "original_content_reviewed": False,
                },
            }
            (page_dir / "visual-gate.json").write_text(json.dumps(gate), encoding="utf-8")

            with self.assertRaisesRegex(SystemExit, "original-content comparison"):
                validate_direct_visual_gate(
                    run_dir, {"page_id": "page_001", "page_index": 1}, page_dir, page_pptx
                )


if __name__ == "__main__":
    unittest.main()
