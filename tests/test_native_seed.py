import tempfile
import unittest
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
SCRIPTS = ROOT / "skills" / "ppt-visual-replicator" / "scripts"

import sys

sys.path.insert(0, str(SCRIPTS))

from seed_native_reconstruction import seed_page  # noqa: E402


class NativeSeedTest(unittest.TestCase):
    def test_seeds_editable_text_shapes_and_images_from_native_pptx(self) -> None:
        with tempfile.TemporaryDirectory() as temp_dir:
            root = Path(temp_dir)
            source = root / "source.pptx"
            image_path = root / "image.png"
            Image.new("RGB", (80, 40), "navy").save(image_path)

            presentation = Presentation()
            slide = presentation.slides.add_slide(presentation.slide_layouts[6])
            text_box = slide.shapes.add_textbox(
                Inches(0.8), Inches(0.5), Inches(5.0), Inches(0.8)
            )
            run = text_box.text_frame.paragraphs[0].add_run()
            run.text = "中文标题"
            run.font.name = "微软雅黑"
            run.font.size = Pt(28)
            native_shape = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(0.8),
                Inches(1.6),
                Inches(4.5),
                Inches(1.0),
            )
            native_shape.fill.solid()
            native_shape.fill.fore_color.rgb = RGBColor(0x00, 0x55, 0x87)
            slide.shapes.add_picture(
                str(image_path), Inches(6.0), Inches(1.5), width=Inches(2.0)
            )
            presentation.save(source)

            loaded = Presentation(source)
            page_dir = root / "page"
            page_dir.mkdir()
            result = seed_page(
                loaded,
                source,
                1,
                page_dir,
                {
                    "source_size_px": {"width": 1280, "height": 720},
                    "slide": {"width": 13.333, "height": 7.5},
                    "content_box": [0, 0, 1280, 720],
                },
            )

            self.assertEqual(result["text_boxes"], 1)
            self.assertGreaterEqual(result["shapes"], 1)
            self.assertEqual(result["images"], 1)
            self.assertEqual(result["unsupported"], 0)
            manifest = (page_dir / "native-manifest-seed.json").read_text(
                encoding="utf-8"
            )
            self.assertIn("中文标题", manifest)
            self.assertIn("Noto Sans CJK SC", manifest)
            self.assertTrue((page_dir / "native-source-slide.pptx").is_file())
            self.assertTrue(any((page_dir / "native-assets").iterdir()))


if __name__ == "__main__":
    unittest.main()
